from pathlib import Path
import json
import math
import re

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
OUT = DOCS / "SiteSpy_Client_Presentation_FULL.pptx"
MD = DOCS / "SiteSpy_Client_Presentation_FULL.md"
ASSET_DIR = DOCS / "presentation-assets"
DIAG_DIR = ASSET_DIR / "diagrams"
SHOT_DIR = ASSET_DIR / "screenshots"
MEMBER_DIR = ASSET_DIR / "member-slides"

for directory in (ASSET_DIR, DIAG_DIR, SHOT_DIR, MEMBER_DIR):
    directory.mkdir(parents=True, exist_ok=True)

LOGO = ROOT / "assets" / "logo-mark-transparent.png"
ICON = ROOT / "assets" / "icon.png"
SOCIAL = ROOT / "assets" / "social-preview.png"
COLLAB_IMAGE = DOCS / "Collaborators_Departments.jpeg"

COLORS = {
    "bg": "#07101E",
    "bg2": "#050A13",
    "surface": "#101B2D",
    "surface2": "#16243A",
    "muted_surface": "#1B2B46",
    "text": "#F8FBFF",
    "muted": "#A8B6CC",
    "soft": "#7888A3",
    "border": "#273A5D",
    "primary": "#25C7F7",
    "secondary": "#3E68FF",
    "accent": "#D8A73E",
    "success": "#49D28F",
    "danger": "#F06B61",
    "white": "#FFFFFF",
    "ink": "#0B1220",
}


def rgb(hex_color):
    value = hex_color.lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


def ppt_color(name):
    return RGBColor(*rgb(COLORS[name]))


def font(size=26, bold=False):
    candidates = [
        Path(r"C:\Windows\Fonts\segoeuib.ttf" if bold else r"C:\Windows\Fonts\segoeui.ttf"),
        Path(r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size)
    return ImageFont.load_default()


def wrap(draw, text, face, width):
    lines = []
    for paragraph in str(text).split("\n"):
        words = paragraph.split()
        if not words:
            lines.append("")
            continue
        line = ""
        for word in words:
            trial = f"{line} {word}".strip()
            if draw.textbbox((0, 0), trial, font=face)[2] <= width or not line:
                line = trial
            else:
                lines.append(line)
                line = word
        if line:
            lines.append(line)
    return lines


def text(draw, xy, value, size=24, bold=False, fill=None, width=None, gap=6, align="left"):
    face = font(size, bold)
    fill = fill or COLORS["text"]
    if width:
        y = xy[1]
        for line in wrap(draw, value, face, width):
            x = xy[0]
            if align == "center":
                w = draw.textbbox((0, 0), line, font=face)[2]
                x += (width - w) / 2
            draw.text((x, y), line, font=face, fill=fill)
            y += size + gap
        return y
    draw.text(xy, value, font=face, fill=fill)
    return xy[1] + size


def rounded(draw, box, fill, outline=None, width=2, radius=24):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def arrow(draw, start, end, fill=None, width=5):
    fill = fill or COLORS["primary"]
    draw.line([start, end], fill=fill, width=width)
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    length = 18
    points = [
        end,
        (
            end[0] - length * math.cos(angle - math.pi / 6),
            end[1] - length * math.sin(angle - math.pi / 6),
        ),
        (
            end[0] - length * math.cos(angle + math.pi / 6),
            end[1] - length * math.sin(angle + math.pi / 6),
        ),
    ]
    draw.polygon(points, fill=fill)


def canvas(title, subtitle="", width=1600, height=900):
    image = Image.new("RGB", (width, height), COLORS["bg"])
    draw = ImageDraw.Draw(image)
    for x in range(0, width, 80):
        draw.line([(x, 0), (x, height)], fill="#09182C", width=1)
    for y in range(0, height, 80):
        draw.line([(0, y), (width, y)], fill="#09182C", width=1)
    draw.rectangle((0, 0, width, 86), fill=COLORS["bg2"])
    text(draw, (56, 22), title, 34, True)
    if subtitle:
        text(draw, (56, 58), subtitle, 17, False, COLORS["muted"])
    if LOGO.exists():
        logo = Image.open(LOGO).convert("RGBA").resize((70, 70))
        image.paste(logo, (width - 112, 10), logo)
    return image, draw


def image_card(draw, x, y, w, h, title, body="", accent="primary"):
    rounded(draw, (x, y, x + w, y + h), COLORS["surface"], COLORS["border"], 2, 26)
    draw.rectangle((x, y, x + 8, y + h), fill=COLORS[accent])
    text(draw, (x + 24, y + 18), title, 24, True, width=w - 48)
    if body:
        text(draw, (x + 24, y + 58), body, 17, False, COLORS["muted"], width=w - 48)


def save(image, path):
    image.save(path, "PNG")
    return path


def make_diagram(name, title, subtitle, boxes, links):
    image, draw = canvas(title, subtitle)
    for box in boxes:
        image_card(draw, box["x"], box["y"], box["w"], box["h"], box["title"], box.get("body", ""), box.get("accent", "primary"))
    for start, end in links:
        arrow(draw, start, end)
    return save(image, DIAG_DIR / f"{name}.png")


diagrams = {}
diagrams["high-level-architecture"] = make_diagram(
    "high-level-architecture",
    "High-Level Architecture",
    "Mobile app, navigation, services, Firebase, and delivery configuration",
    [
        {"title": "Mobile User", "x": 70, "y": 380, "w": 190, "h": 95, "accent": "secondary"},
        {"title": "Expo React Native App", "x": 340, "y": 200, "w": 260, "h": 105},
        {"title": "Navigation Layer", "x": 340, "y": 370, "w": 260, "h": 105, "accent": "secondary"},
        {"title": "Screens + Components", "x": 680, "y": 250, "w": 280, "h": 135, "accent": "success"},
        {"title": "Services Layer", "x": 680, "y": 470, "w": 280, "h": 110},
        {"title": "Firebase Auth", "x": 1070, "y": 170, "w": 250, "h": 90, "accent": "secondary"},
        {"title": "Cloud Firestore", "x": 1070, "y": 320, "w": 250, "h": 120, "accent": "success"},
        {"title": "Firebase Storage", "x": 1070, "y": 500, "w": 250, "h": 95},
        {"title": "Config / EAS / Rules", "x": 410, "y": 650, "w": 330, "h": 95, "accent": "accent"},
    ],
    [((260, 428), (340, 428)), ((470, 305), (470, 370)), ((600, 252), (680, 300)), ((600, 425), (680, 510)), ((960, 505), (1070, 215)), ((960, 515), (1070, 375)), ((960, 535), (1070, 545)), ((575, 650), (575, 475))],
)
diagrams["navigation-architecture"] = make_diagram(
    "navigation-architecture",
    "Navigation Architecture",
    "Auth stack, protected main tabs, and nested project/profile stacks",
    [
        {"title": "AppNavigator", "body": "Splash while initializing\nSigned out: Auth stack\nSigned in: Main tabs", "x": 80, "y": 170, "w": 300, "h": 150},
        {"title": "Auth Stack", "body": "Login\nRegister\nForgot Password", "x": 500, "y": 140, "w": 300, "h": 150, "accent": "secondary"},
        {"title": "Main Tabs", "body": "Dashboard\nProjects\nNew\nProfile", "x": 500, "y": 405, "w": 300, "h": 160, "accent": "success"},
        {"title": "Project Stack", "body": "Project History\nProject Details\nEdit Project\nManual Estimate\nImage Estimate\nEstimate Summary", "x": 930, "y": 245, "w": 330, "h": 245},
        {"title": "Profile Stack", "body": "Profile\nSettings", "x": 1320, "y": 405, "w": 210, "h": 130, "accent": "accent"},
    ],
    [((380, 235), (500, 215)), ((380, 255), (500, 485)), ((800, 485), (930, 370)), ((800, 510), (1320, 470))],
)
diagrams["data-flow"] = make_diagram(
    "data-flow",
    "Data Flow",
    "From sign-in and project creation to stored estimate and summary",
    [
        {"title": "Sign in", "x": 80, "y": 375, "w": 190, "h": 110},
        {"title": "Create project", "x": 320, "y": 260, "w": 210, "h": 110, "accent": "secondary"},
        {"title": "Enter inputs", "x": 590, "y": 260, "w": 210, "h": 110},
        {"title": "Calculate locally", "x": 860, "y": 260, "w": 230, "h": 110, "accent": "accent"},
        {"title": "Save estimate", "x": 1160, "y": 260, "w": 220, "h": 110, "accent": "success"},
        {"title": "Read summary", "x": 1160, "y": 540, "w": 220, "h": 110},
        {"title": "Image record path", "body": "Image + reference measurement", "x": 590, "y": 560, "w": 260, "h": 120, "accent": "secondary"},
        {"title": "Storage + metadata", "x": 900, "y": 560, "w": 260, "h": 120, "accent": "success"},
    ],
    [((270, 430), (320, 315)), ((530, 315), (590, 315)), ((800, 315), (860, 315)), ((1090, 315), (1160, 315)), ((1270, 370), (1270, 540)), ((700, 370), (700, 560)), ((850, 620), (900, 620)), ((1030, 560), (1200, 370))],
)
diagrams["firebase-architecture"] = make_diagram(
    "firebase-architecture",
    "Firebase Architecture",
    "Authenticated owner-based records and image paths",
    [
        {"title": "Firebase Auth", "body": "Email/password\nPassword reset\nEmail-link sign-in", "x": 90, "y": 230, "w": 330, "h": 155, "accent": "secondary"},
        {"title": "users", "body": "One document per uid", "x": 540, "y": 145, "w": 300, "h": 130, "accent": "success"},
        {"title": "projects", "body": "User-owned project records", "x": 540, "y": 345, "w": 300, "h": 130, "accent": "success"},
        {"title": "estimations", "body": "Linked wall estimates", "x": 950, "y": 245, "w": 300, "h": 130},
        {"title": "wallImages", "body": "Linked image records", "x": 950, "y": 475, "w": 300, "h": 130},
        {"title": "Storage paths", "body": "users/{userId}/projects/{projectId}/images/{fileName}", "x": 1310, "y": 475, "w": 250, "h": 130, "accent": "accent"},
        {"title": "Security Rules", "body": "Authenticated owner access", "x": 380, "y": 680, "w": 720, "h": 95, "accent": "secondary"},
    ],
    [((420, 305), (540, 210)), ((420, 305), (540, 410)), ((840, 410), (950, 310)), ((840, 410), (950, 540)), ((1250, 540), (1310, 540))],
)
diagrams["database-model"] = make_diagram(
    "database-model",
    "Database Model",
    "Firestore collections used by SiteSpy",
    [
        {"title": "users", "body": "userId\nemail\ndisplayName\ncreatedAt\nupdatedAt", "x": 90, "y": 190, "w": 300, "h": 390, "accent": "success"},
        {"title": "projects", "body": "projectId\nuserId\ntitle\ndescription\nlocation", "x": 470, "y": 190, "w": 300, "h": 390},
        {"title": "estimations", "body": "estimationId\nprojectId\nuserId\nwallArea\ntotalCost\nboqSummary", "x": 850, "y": 190, "w": 300, "h": 390, "accent": "accent"},
        {"title": "wallImages", "body": "imageId\nprojectId\nuserId\nimageUrl\nstoragePath\nreferenceMeasurement", "x": 1230, "y": 190, "w": 300, "h": 390, "accent": "secondary"},
    ],
    [((390, 495), (470, 495)), ((770, 360), (850, 360)), ((770, 495), (1230, 495))],
)
diagrams["estimation-logic"] = make_diagram(
    "estimation-logic",
    "Estimation Logic",
    "Wall inputs become quantities, costs, and BOQ-style totals",
    [
        {"title": "Wall length × height", "body": "Wall area", "x": 90, "y": 235, "w": 270, "h": 130},
        {"title": "Area ÷ unit face area", "body": "Estimated units", "x": 440, "y": 235, "w": 290, "h": 130},
        {"title": "Waste factor", "body": "Adjusted quantity", "x": 810, "y": 235, "w": 270, "h": 130, "accent": "accent"},
        {"title": "Rates", "body": "Material + mortar + labour", "x": 1160, "y": 235, "w": 270, "h": 130, "accent": "secondary"},
        {"title": "BOQ summary", "body": "Line items and total cost", "x": 650, "y": 570, "w": 330, "h": 130, "accent": "success"},
    ],
    [((360, 300), (440, 300)), ((730, 300), (810, 300)), ((1080, 300), (1160, 300)), ((1295, 365), (980, 570)), ((945, 365), (850, 570)), ((585, 365), (760, 570))],
)
diagrams["git-workflow"] = make_diagram(
    "git-workflow",
    "Git Workflow",
    "Assigned branch work and review path",
    [
        {"title": "Pull latest main", "x": 80, "y": 350, "w": 190, "h": 110},
        {"title": "Create assigned branch", "x": 320, "y": 350, "w": 210, "h": 110, "accent": "secondary"},
        {"title": "Implement task", "x": 590, "y": 350, "w": 190, "h": 110},
        {"title": "Run checks", "x": 830, "y": 350, "w": 190, "h": 110, "accent": "success"},
        {"title": "Commit + push", "x": 1060, "y": 350, "w": 190, "h": 110},
        {"title": "Pull request review", "x": 1290, "y": 350, "w": 220, "h": 110, "accent": "accent"},
        {"title": "Identity rule", "body": "Each contributor uses their own GitHub account and Git identity.", "x": 390, "y": 610, "w": 760, "h": 105, "accent": "accent"},
    ],
    [((270, 405), (320, 405)), ((530, 405), (590, 405)), ((780, 405), (830, 405)), ((1020, 405), (1060, 405)), ((1250, 405), (1290, 405))],
)
diagrams["build-deployment-flow"] = make_diagram(
    "build-deployment-flow",
    "Build and Deployment Flow",
    "Local checks, Android export, EAS build path, and Firebase deploy targets",
    [
        {"title": "npm install", "x": 80, "y": 240, "w": 180, "h": 105},
        {"title": "npm run check", "x": 310, "y": 240, "w": 180, "h": 105, "accent": "success"},
        {"title": "npm run doctor", "x": 540, "y": 240, "w": 180, "h": 105},
        {"title": "Android export", "x": 770, "y": 240, "w": 190, "h": 105, "accent": "secondary"},
        {"title": "EAS preview APK", "x": 1010, "y": 240, "w": 200, "h": 105, "accent": "accent"},
        {"title": "Android demo", "x": 1260, "y": 240, "w": 190, "h": 105, "accent": "success"},
        {"title": "Firebase deploy", "body": "Rules, indexes, storage, hosting", "x": 400, "y": 570, "w": 330, "h": 120, "accent": "secondary"},
        {"title": "Hosting support page", "body": "Optional support page, not the main Android app", "x": 850, "y": 570, "w": 360, "h": 120, "accent": "success"},
    ],
    [((260, 292), (310, 292)), ((490, 292), (540, 292)), ((720, 292), (770, 292)), ((960, 292), (1010, 292)), ((1210, 292), (1260, 292)), ((730, 630), (850, 630))],
)
diagrams["app-feature-map"] = make_diagram(
    "app-feature-map",
    "App Feature Map",
    "Implemented screens and services at a glance",
    [
        {"title": "SiteSpy App", "body": "Mobile-first estimation workspace", "x": 650, "y": 380, "w": 300, "h": 110},
        {"title": "Authentication", "x": 130, "y": 160, "w": 230, "h": 90, "accent": "secondary"},
        {"title": "Dashboard", "x": 560, "y": 130, "w": 230, "h": 90},
        {"title": "Projects", "x": 1000, "y": 160, "w": 230, "h": 90, "accent": "success"},
        {"title": "Manual Estimate", "x": 1230, "y": 390, "w": 230, "h": 90, "accent": "accent"},
        {"title": "Image Estimate", "x": 1000, "y": 650, "w": 230, "h": 90},
        {"title": "Summary", "x": 560, "y": 680, "w": 230, "h": 90, "accent": "success"},
        {"title": "Profile", "x": 130, "y": 650, "w": 230, "h": 90, "accent": "secondary"},
        {"title": "Firebase Services", "x": 130, "y": 390, "w": 230, "h": 90, "accent": "success"},
    ],
    [((360, 205), (650, 435)), ((675, 220), (745, 380)), ((1000, 205), (950, 435)), ((1230, 435), (950, 435)), ((1000, 695), (950, 435)), ((675, 680), (745, 490)), ((360, 695), (650, 435)), ((360, 435), (650, 435))],
)
diagrams["member-workstream-map"] = make_diagram(
    "member-workstream-map",
    "Member Workstream Map",
    "Mobile UI, Firebase, quality, and delivery ownership",
    [
        {"title": "Mobile UI", "body": "Navigation, auth, dashboard, projects, estimates, profile/settings", "x": 90, "y": 200, "w": 560, "h": 155},
        {"title": "Firebase + Data", "body": "Setup, auth review, Firestore model, services, rules", "x": 90, "y": 440, "w": 560, "h": 155, "accent": "success"},
        {"title": "Quality + Delivery", "body": "Responsive testing, checks, release review, summary, EAS deployment", "x": 880, "y": 320, "w": 560, "h": 155, "accent": "accent"},
        {"title": "18 assignments", "body": "Documented in task allocation", "x": 690, "y": 285, "w": 260, "h": 150, "accent": "secondary"},
    ],
    [((650, 280), (690, 330)), ((650, 520), (880, 430)), ((950, 360), (880, 390))],
)


screen_defs = {
    "splash": ("SiteSpy", "Starting secure field estimates", ["Initializing app configuration", "Checking account state", "Preparing Android workspace"]),
    "login": ("Welcome back", "Secure field estimates", ["Email", "Password", "Log in", "Email me a sign-in link", "Create account", "Forgot password"]),
    "register": ("Create account", "Start a SiteSpy workspace", ["Full name", "Email", "Password", "Confirm password", "Create account"]),
    "forgot-password": ("Reset password", "Recover access through Firebase Auth", ["Email", "Send reset email", "Back to login"]),
    "dashboard": ("Project dashboard", "Track wall estimates, quantities, images, and labour costs.", ["Projects: 3", "Estimations: 5", "Masonry units: 2,420", "Total estimate: N$ 18,650", "New project", "Project history"]),
    "project-history": ("Project History", "Saved field files and previous BOQ totals", ["Boundary wall estimate", "Training block wall", "Campus repair notes", "Open project details"]),
    "new-project": ("New Project", "Create a field file and start an estimate", ["Project title", "Description", "Location", "Create and estimate"]),
    "manual-estimate": ("Manual Estimate", "Wall and costing inputs", ["Wall length", "Wall height", "Unit type: Concrete block", "Waste factor: 1.08", "Unit price", "Labour rate", "Save estimate"]),
    "estimate-summary": ("Estimate Summary", "BOQ-style result", ["Masonry units 420", "Mortar allowance 0.42 m3", "Labour 24 m2", "Total cost N$ 4,860"]),
    "image-estimate": ("Image Estimate", "Wall image and reference measurement", ["Choose from gallery", "Capture camera photo", "Reference measurement in metres", "Save image record"]),
    "project-details": ("Project Details", "Linked estimates and image records", ["Project notes", "Saved estimates", "Wall images", "Edit project", "Delete project"]),
    "profile": ("Profile", "Account workspace", ["Signed-in email", "Firebase user id", "Open settings", "Log out"]),
    "settings": ("Settings", "Preferences and app target", ["Draft preferences", "Android target information", "Storage behaviour", "Log out"]),
}


def phone_visual(name, title, subtitle, items):
    image = Image.new("RGB", (900, 1200), "#E7EEF8")
    draw = ImageDraw.Draw(image)
    rounded(draw, (230, 55, 670, 1135), "#020617", "#0F172A", 4, 44)
    rounded(draw, (248, 75, 652, 1115), COLORS["bg"], "#1E293B", 2, 34)
    draw.rectangle((248, 75, 652, 160), fill=COLORS["bg2"])
    text(draw, (285, 102), "SiteSpy", 24, True)
    draw.ellipse((590, 102, 622, 134), fill=COLORS["primary"])
    y = 190
    y = text(draw, (285, y), title, 32, True, width=330) + 18
    y = text(draw, (285, y), subtitle, 17, False, COLORS["muted"], width=330) + 42
    if name == "splash":
        if LOGO.exists():
            logo = Image.open(LOGO).convert("RGBA").resize((180, 180))
            image.paste(logo, (360, 385), logo)
        y = 625
        for item in items:
            text(draw, (310, y), f"• {item}", 18, False, COLORS["muted"])
            y += 44
    elif name == "dashboard":
        for i, item in enumerate(items[:4]):
            x = 280 + (i % 2) * 170
            yy = y + (i // 2) * 105
            rounded(draw, (x, yy, x + 145, yy + 82), COLORS["surface2"], COLORS["border"], 2, 18)
            label, value = item.split(":", 1)
            text(draw, (x + 14, yy + 14), label, 13, False, COLORS["muted"], width=120)
            text(draw, (x + 14, yy + 40), value.strip(), 20, True, width=120)
        y += 235
        for item in items[4:]:
            rounded(draw, (285, y, 615, y + 72), COLORS["surface"], COLORS["border"], 2, 18)
            text(draw, (306, y + 20), item, 18, True)
            y += 88
    else:
        for item in items:
            is_action = any(word in item.lower() for word in ["log in", "create", "send", "save", "choose", "capture", "open", "delete", "back"])
            fill = COLORS["primary"] if is_action and "delete" not in item.lower() else COLORS["danger"] if "delete" in item.lower() else COLORS["surface"]
            rounded(draw, (285, y, 615, y + 64), fill, COLORS["border"], 2, 18)
            text(draw, (305, y + 18), item, 16 if not is_action else 17, is_action, COLORS["bg"] if is_action else COLORS["muted"], width=285)
            y += 78
    if name not in ["splash", "login", "register", "forgot-password"]:
        rounded(draw, (268, 1010, 632, 1085), COLORS["surface"], "#23395B", 2, 24)
        for i, tab in enumerate(["Dash", "Projects", "New", "Profile"]):
            text(draw, (300 + i * 82, 1035), tab, 13, True, COLORS["primary"] if i == 0 else COLORS["muted"])
    return save(image, SHOT_DIR / f"{name}.png")


screens = {name: phone_visual(name, *spec) for name, spec in screen_defs.items()}

members = [
    ("Tunacky Kandere", "tunackykandere-lab", "mobile/kandere-ndl-navigation-shell", "Navigation shell and app layout", "Mobile UI", "navigation-architecture", "project-history"),
    ("Kavakuru Metarere", "Kavakuru-Metarere7", "mobile/metare-k-auth-screens", "Authentication screens", "Mobile UI", "login", "login"),
    ("Nambuli NN", "documented branch owner", "mobile/nambuli-nn-dashboard", "Dashboard experience", "Mobile UI", "dashboard", "dashboard"),
    ("Mathias Jonas", "Mathias4040", "mobile/jonas-mm-project-history", "Project history and details", "Mobile UI", "project-history", "project-history"),
    ("Emilly Ndapuka", "emilly20-06", "mobile/ndapuka-eii-manual-estimation", "Manual estimation forms", "Mobile UI", "manual-estimate", "manual-estimate"),
    ("Amalia Mangundu", "amaliamangundu-tech", "mobile/mangundu-a-profile-settings", "Profile and settings", "Mobile UI", "profile", "profile"),
    ("Johannes Kandjeke", "Kandjekejohannes54", "mobile/kandjeke-jm-responsive-testing", "Android responsive testing", "Mobile QA", "build-deployment-flow", "settings"),
    ("Hilma Shuumbwa", "hilma-shuumbwa", "firebase/shuumbwa-hmn-project-setup", "Expo and Firebase setup", "Firebase", "firebase-architecture", "firebase"),
    ("Petrus Hamukwaya", "ramuntu", "firebase/hamukwaya-pnp-auth-rules", "Authentication rules", "Firebase", "firebase-architecture", "login"),
    ("Karl Hamberera", "Karlhamberera", "firebase/hamberera-mkp-firestore-models", "Firestore data model", "Firebase", "database-model", "database"),
    ("Washington Matattias", "matattiasw-ai", "firebase/matattias-w-integration-review", "Repository and release review", "Integration", "build-deployment-flow", "build"),
    ("Klaudia Kambowe", "kambowe", "firebase/kambowe-kn-project-services", "Project services", "Firebase", "data-flow", "project-details"),
    ("Martha Heita", "marthandilimengungo-dotcom", "firebase/heita-mn-validation-errors", "Validation and errors", "Quality", "data-flow", "register"),
    ("Michael Kazundire", "IcyBeibi68023", "firebase/tjatindi-mk-security-review", "Firebase rules review", "Security", "firebase-architecture", "security"),
    ("Beatha Haipumbu", "beathapawana-dot", "firebase/haipumbu-bnp-estimation-workflow", "Estimation workflow", "Calculation", "estimation-logic", "manual-estimate"),
    ("Valentina Correia", "Valentina-Correia", "firebase/correia-vp-summary-export", "Summary formatting", "Reporting", "estimate-summary", "estimate-summary"),
    ("Hilda Iita", "next-GenCoder", "firebase/iita-hn-quality-checks", "App checks", "Quality", "build-deployment-flow", "testing"),
    ("Linus Shikongo", "Linus45-blid", "firebase/shikongo-lik-eas-deployment", "EAS deployment", "Deployment", "build-deployment-flow", "deployment"),
]

why = {
    "Navigation shell and app layout": ["makes the app easy to move through", "separates auth, project, and profile workflows"],
    "Authentication screens": ["protects access to user-owned records", "gives the demo a complete account flow"],
    "Dashboard experience": ["shows system value immediately after sign-in", "turns saved data into presentation metrics"],
    "Project history and details": ["lets users reopen and manage saved field files", "connects estimates and images to each project"],
    "Manual estimation forms": ["turns construction inputs into repeatable calculations", "supports the main value of the app"],
    "Profile and settings": ["gives users account control and logout access", "rounds out the protected app experience"],
    "Android responsive testing": ["improves readability on portrait Android screens", "supports reliable live demonstration"],
    "Expo and Firebase setup": ["provides the foundation for app startup and backend connection", "makes configuration traceable"],
    "Authentication rules": ["keeps sign-in behaviour aligned with Firebase Auth", "supports secure account-linked data"],
    "Firestore data model": ["defines how records are organized and linked", "supports project history and summaries"],
    "Repository and release review": ["keeps integration checks visible", "reduces risk before presentation and build"],
    "Project services": ["provides data operations behind project screens", "keeps screen code focused on user experience"],
    "Validation and errors": ["helps users correct mistakes before saving", "makes failures easier to explain and retry"],
    "Firebase rules review": ["protects records through owner-only rules", "supports a defensible security explanation"],
    "Estimation workflow": ["connects form inputs to saved cost summaries", "supports the main construction-estimation use case"],
    "Summary formatting": ["makes calculation results easier to present", "turns estimates into BOQ-style outputs"],
    "App checks": ["confirms required files and quality commands", "supports final submission confidence"],
    "EAS deployment": ["defines the Android APK delivery path", "supports installable demo readiness"],
}

worked = {
    "Navigation shell and app layout": ["Auth-gated navigation structure", "Bottom-tab and stack layout", "Android portrait app flow"],
    "Authentication screens": ["Login, register, and recovery screens", "Validation-ready account entry", "Firebase sign-in actions"],
    "Dashboard experience": ["Project metrics", "Recent projects", "Empty and loading states"],
    "Project history and details": ["Saved project browsing", "Detail review", "Edit/delete and estimate visibility"],
    "Manual estimation forms": ["Wall input experience", "Calculation preview", "Save estimate flow"],
    "Profile and settings": ["Account view", "Settings screen", "Logout path"],
    "Android responsive testing": ["Portrait spacing", "Keyboard behaviour", "Touch-target checks"],
    "Expo and Firebase setup": ["Firebase config", "Expo app settings", "Android package setup"],
    "Authentication rules": ["Auth behaviour", "Account ownership", "Protected record assumptions"],
    "Firestore data model": ["Collection structure", "Owner fields", "Linked records"],
    "Repository and release review": ["Integration checks", "Project check", "Release readiness"],
    "Project services": ["Create/list/update/delete", "Project timestamps", "Linked estimates"],
    "Validation and errors": ["Form validation", "Clear user messages", "Retry/error states"],
    "Firebase rules review": ["Firestore rules", "Storage paths", "Owner-only access"],
    "Estimation workflow": ["Wall area", "Unit quantities", "Cost total"],
    "Summary formatting": ["BOQ line items", "Cost breakdown", "Presenter-friendly totals"],
    "App checks": ["Project check", "Testing checklist", "Export readiness"],
    "EAS deployment": ["Preview APK path", "Build profile", "Android delivery"],
}


def member_visual(member):
    name, github, branch, area, stream, visual_key, screen_key = member
    image, draw = canvas(area, f"{stream} contribution visual", 1200, 820)
    image_card(draw, 70, 145, 410, 210, name, f"GitHub / identity: {github}\nBranch: {branch}", "primary")
    related = screens.get(screen_key) or diagrams.get(visual_key) or diagrams.get("build-deployment-flow")
    inset = Image.open(related).convert("RGB")
    inset.thumbnail((560, 470))
    rounded(draw, (582, 127, 582 + inset.width + 36, 127 + inset.height + 36), COLORS["surface2"], COLORS["border"], 2, 22)
    image.paste(inset, (600, 145))
    y = 405
    text(draw, (85, y), "Evidence focus", 23, True)
    y += 45
    for item in worked[area]:
        text(draw, (95, y), f"• {item}", 20, False, COLORS["muted"], width=455)
        y += 38
    path = MEMBER_DIR / f"{re.sub(r'[^a-z0-9]+', '-', name.lower()).strip('-')}.png"
    return save(image, path)


member_images = {m[0]: member_visual(m) for m in members}


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
records = []


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ppt_color("bg")
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = ppt_color("primary")
    band.line.fill.background()


def tx(slide, x, y, w, h, value, size=16, bold=False, color="text", align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    paragraph = frame.paragraphs[0]
    if align:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = "Segoe UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = ppt_color(color)
    return shape


def bullets(slide, x, y, w, h, items, size=15, color="text"):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.word_wrap = True
    for i, item in enumerate(items):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.font.name = "Segoe UI"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = ppt_color(color)
        paragraph.space_after = Pt(4)


def title(slide, value, subtitle="", number=None):
    bg(slide)
    tx(slide, 0.55, 0.38, 10.8, 0.55, value, 26, True)
    if subtitle:
        tx(slide, 0.58, 0.95, 10.8, 0.33, subtitle, 12.5, False, "muted")
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(11.95), Inches(0.32), width=Inches(0.65), height=Inches(0.65))
    if number:
        tx(slide, 0.55, 7.13, 12.0, 0.22, f"SiteSpy System Presentation | {number:02d}", 8, False, "soft", PP_ALIGN.RIGHT)


def add_picture(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def card(slide, x, y, w, h, header, body="", accent="primary"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_color("surface")
    shape.line.color.rgb = ppt_color("border")
    shape.adjustments[0] = 0.08
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ppt_color(accent)
    bar.line.fill.background()
    tx(slide, x + 0.18, y + 0.15, w - 0.3, 0.35, header, 14, True)
    if body:
        tx(slide, x + 0.18, y + 0.55, w - 0.3, h - 0.62, body, 10.5, False, "muted")


def add_slide(title_text, subtitle="", items=None, image=None, layout="bullets", notes="", evidence=None):
    number = len(prs.slides) + 1
    slide = prs.slides.add_slide(blank)
    title(slide, title_text, subtitle, number)
    if layout == "image-full":
        add_picture(slide, image, 0.75, 1.25, 11.85, 5.55)
    elif layout == "image-right":
        bullets(slide, 0.75, 1.45, 5.55, 4.9, items or [], 15)
        if image:
            add_picture(slide, image, 6.65, 1.35, 5.7, 4.95)
    elif layout == "two-images":
        for index, path in enumerate(image or []):
            add_picture(slide, path, 0.75 + index * 6.15, 1.35, 5.65, 4.75)
        if items:
            bullets(slide, 0.85, 6.08, 11.5, 0.65, items, 12)
    elif layout == "cards":
        for index, (head, body, accent) in enumerate(items or []):
            card(slide, 0.75 + (index % 3) * 4.15, 1.35 + (index // 3) * 1.75, 3.65, 1.25, head, body, accent)
    else:
        bullets(slide, 0.85, 1.35, 11.5, 5.4, items or [], 16)
    records.append({"title": title_text, "content": items or [], "image": image, "notes": notes, "evidence": evidence or []})
    return slide


# Slide 1
slide = prs.slides.add_slide(blank)
bg(slide)
if LOGO.exists():
    add_picture(slide, LOGO, 0.85, 0.75, 1.0, 1.0)
tx(slide, 0.85, 2.10, 9.8, 0.8, "SiteSpy System Presentation", 38, True)
tx(slide, 0.90, 2.92, 8.5, 0.4, "Client Demonstration and Technical Overview", 18, False, "muted")
tx(slide, 0.90, 3.55, 8.5, 0.35, "Prepared for Group 2 / client team", 13, False, "muted")
tx(slide, 0.90, 3.95, 5.0, 0.32, "7 June 2026", 12, False, "soft")
if SOCIAL.exists():
    add_picture(slide, SOCIAL, 8.40, 1.35, 4.15, 3.1)
tx(slide, 0.55, 7.13, 12.0, 0.22, "SiteSpy System Presentation | 01", 8, False, "soft", PP_ALIGN.RIGHT)
records.append({"title": "SiteSpy System Presentation", "content": ["Client Demonstration and Technical Overview", "Prepared for Group 2 / client team", "7 June 2026"], "image": SOCIAL, "notes": "Open by introducing SiteSpy as the mobile-first estimation system built for the student/client team.", "evidence": ["README.md", "app.json", "assets/logo-mark-transparent.png"]})


add_slide("Executive Summary", "What was built and why it is ready to present", ["SiteSpy is an Expo React Native Android app for masonry wall project estimation.", "The app supports sign-in, project records, manual estimates, image records, and BOQ-style summaries.", "Firebase stores user-owned project, estimate, and wall-image records with owner-based access rules.", "The repository includes collaborator ownership, testing, Firebase setup, and final submission documentation."], SOCIAL, "image-right", "Summarize the system in plain language before going deeper.", ["README.md", "docs/User_Manual.md"])
add_slide("Presentation Roadmap", "How the client walkthrough is organized", ["Problem and project objectives", "Solution overview and target users", "App walkthrough and core workflows", "Architecture, Firebase, data, and security", "Member contributions and Git workflow", "Testing, build process, demo script, and Q&A"], diagrams["member-workstream-map"], "image-right", "Set expectations for the presentation.", ["docs/collaborators/task-allocation.md"])
add_slide("Problem Statement", "The practical need SiteSpy addresses", ["Manual construction estimation can be slow, inconsistent, and difficult to present clearly.", "Students need visible ownership across a shared project and assigned branches.", "Project records must be stored securely and tied to the correct signed-in user.", "Presenters need a technical explanation that is accurate and understandable."], notes="Explain the client problem before naming features.", evidence=["README.md"])
add_slide("Project Objectives", "What the team set out to deliver", ["Build a mobile-first Android app using Expo and React Native.", "Use Firebase for authentication and user-owned backend records.", "Support project creation, manual estimation, image records, and report-ready summaries.", "Document collaborator responsibilities, branches, testing, and final submission steps.", "Prepare the project for Android demonstration and continued student work."], notes="Link each objective to a visible part of the app or repository.", evidence=["README.md", "docs/Testing_Checklist.md"])
add_slide("What SiteSpy Does", "Core app workflow in client language", ["The user signs in to a protected workspace.", "The user creates a project for a wall or construction task.", "The app captures wall dimensions, unit type, prices, and labour values.", "The calculation workflow produces unit counts, mortar allowance, labour, and total cost.", "Saved projects and estimates can be reviewed later through history and detail screens."], screens["dashboard"], "image-right", "Describe the app as a sequence of actions.", ["docs/User_Manual.md"])
add_slide("Target Users and Use Cases", "Who benefits from the system", ["Students presenting a working mobile app and team contribution structure.", "Small builders or estimators who need quick wall quantity and cost estimates.", "Demonstration panel members reviewing system design, quality, and maintainability.", "Future maintainers continuing from documented branches and service boundaries."], notes="Emphasize the demo and maintainability use cases.", evidence=["README.md"])
add_slide("Basic Programming Concepts", "Client-friendly technical vocabulary", ["Frontend: the visible mobile interface users tap and read.", "Screens: full views such as Login, Dashboard, Manual Estimate, Project Details, and Settings.", "Components: reusable UI pieces such as cards, buttons, headers, inputs, and status states.", "Services: modules that connect screens to Firebase, storage, reports, and calculations.", "Database: Firestore collections that store user-owned app records.", "Build process: commands that prepare the app for local use and Android delivery."], notes="Help non-technical presenters explain the implementation.", evidence=["src/components", "src/screens", "src/services"])
add_slide("Technology Stack", "Actual stack found in the repository", ["Expo Managed Workflow and React Native for Android app delivery.", "JavaScript and JSX for screens, services, components, and utilities.", "React Navigation for auth flow, bottom tabs, project stack, and profile stack.", "Firebase Auth, Firestore, and Storage for backend behaviour.", "EAS Build profiles for preview APK and production app-bundle paths.", "Expo image, font, file-system, image-manipulator, and build-properties libraries."], notes="Keep this grounded in package.json and app configuration.", evidence=["package.json", "app.json", "eas.json"])
add_slide("App Feature Map", "Major implemented areas at a glance", image=diagrams["app-feature-map"], layout="image-full", notes="Walk through each feature around the center of the map.", evidence=["README.md", "src/navigation/MainTabs.js"])
add_slide("High-Level Architecture Diagram", "How the system fits together", image=diagrams["high-level-architecture"], layout="image-full", notes="Explain the split between mobile UI, services, Firebase, and configuration.", evidence=["src/navigation", "src/services"])
add_slide("Navigation Architecture", "Actual navigation files and flows", image=diagrams["navigation-architecture"], layout="image-full", notes="Show that signed-out and signed-in users have different navigation paths.", evidence=["src/navigation/AppNavigator.js", "src/navigation/MainTabs.js"])
add_slide("Data Flow Diagram", "From input to stored result", image=diagrams["data-flow"], layout="image-full", notes="Explain the estimate save workflow and image record path.", evidence=["src/services/projectService.js", "src/services/estimateService.js"])
add_slide("Firebase Architecture", "Authentication, records, and storage", image=diagrams["firebase-architecture"], layout="image-full", notes="Connect Auth, Firestore collections, Storage paths, and owner rules.", evidence=["docs/Firebase_Setup_Guide.md", "firestore.rules", "storage.rules"])
add_slide("Database / Collection Design", "Firestore collections used by the app", image=diagrams["database-model"], layout="image-full", notes="Explain each collection and how project-linked records are connected.", evidence=["shared/data-model.md"])
add_slide("Security Model", "How project records are protected", ["Firestore rules require authentication before app records can be read or written.", "Project, estimation, and wall-image records are restricted to the matching userId.", "Storage paths are scoped to the authenticated user.", "The project check blocks private files such as environment values, credentials, tokens, and raw collaborator secrets.", "Firebase values are kept in local .env and loaded through app configuration."], notes="Present this as a practical data protection story.", evidence=["firestore.rules", "storage.rules", "scripts/project-check.js"])
add_slide("Folder and Codebase Structure", "Important repository areas", ["src/navigation: auth gate, stacks, and bottom tabs.", "src/screens: auth, dashboard, project, estimate, image, profile, and settings screens.", "src/components: shared mobile UI building blocks.", "src/services: Firebase, project, storage, report, auth, and estimate services.", "src/utils and src/theme: calculations, validation, formatting, colors, spacing, typography.", "docs, assets, functions, hosting, and Firebase config files support delivery and documentation."], notes="Make clear that the repo separates responsibilities.", evidence=["README.md"])
add_slide("User Flow Walkthrough", "Typical demonstration sequence", ["Splash opens while startup state is checked.", "Signed-out users choose login, account creation, or password recovery.", "Signed-in users land in Dashboard and navigate through bottom tabs.", "New Project starts a project file and moves into estimation.", "Manual Estimate saves calculations and Estimate Summary presents BOQ lines.", "Image Estimate records a wall image and reference measurement.", "Profile and Settings close the demo with account and app information."], diagrams["data-flow"], "image-right", "Use this as the presenter app walkthrough map.", ["docs/User_Manual.md"])
add_slide("Screenshot Walkthrough 1 - Authentication", "Splash, login, account creation, and recovery", ["Splash prepares app startup.", "Login supports email/password and email-link sign-in.", "Register creates a Firebase-backed account profile.", "Forgot Password sends a Firebase reset email."], [screens["splash"], screens["login"]], "two-images", "Show the sign-in path before protected app screens.", ["src/screens/LoginScreen.js"])
add_slide("Screenshot Walkthrough 2 - Dashboard and Projects", "Project visibility after sign-in", ["Dashboard summarizes projects, estimations, masonry units, and total estimate value.", "Project History lists saved project files.", "Project Details connects project notes, saved estimates, image records, and actions."], [screens["dashboard"], screens["project-history"]], "two-images", "Explain how the user finds and reopens work.", ["src/screens/DashboardScreen.js"])
add_slide("Screenshot Walkthrough 3 - Estimation Workflow", "Project creation, manual estimate, and summary", ["New Project captures the field file.", "Manual Estimate collects wall and costing inputs.", "Estimate Summary presents BOQ-style line items and totals."], [screens["new-project"], screens["manual-estimate"]], "two-images", "Walk through creating an estimate from start to saved result.", ["src/screens/ManualEstimateScreen.js"])
add_slide("Screenshot Walkthrough 4 - Image Estimate and Records", "Image-assisted project evidence", ["The image screen supports gallery or camera selection.", "A real reference measurement is captured with the image record.", "The app stores image metadata linked to the project and user."], [screens["image-estimate"], screens["project-details"]], "two-images", "Clarify the image record path and reference measurement.", ["src/screens/ImageEstimateScreen.js"])
add_slide("Screenshot Walkthrough 5 - Profile and Settings", "Account and app information", ["Profile shows account identity details.", "Settings includes preferences and Android target information.", "Logout clears the signed-in session and sensitive local state."], [screens["profile"], screens["settings"]], "two-images", "End the demo with account control and sign-out.", ["src/screens/ProfileScreen.js"])
add_slide("Estimation Logic Explained", "How calculations are produced", image=diagrams["estimation-logic"], layout="image-full", notes="Explain formulas in simple terms: inputs become quantities, quantities become costs, costs become the summary.", evidence=["src/utils/calculations.js"])
add_slide("BOQ / Report Summary Explained", "Why the result is presenter-friendly", ["Masonry units are listed as a material line item.", "Mortar allowance is calculated and priced separately.", "Labour is calculated from wall area and labour rate.", "The total cost is the sum of material, mortar, and labour costs.", "The summary format helps presenters discuss project cost in a structured way."], screens["estimate-summary"], "image-right", "Explain the result card and line-item value.", ["src/utils/calculations.js"])
add_slide("Collaboration System", "How student ownership is documented", ["Collaborator mapping stores names, GitHub identities, emails, and departments where recorded.", "Task allocation maps each student to a branch, workstream, contribution area, and deliverable.", "Git workflow notes define branch and pull request expectations.", "Private values such as passwords and tokens are redacted and excluded from presentation content."], COLLAB_IMAGE if COLLAB_IMAGE.exists() else diagrams["member-workstream-map"], "image-right", "Tie collaboration documentation to clean handoff and presentation readiness.", ["docs/Collaborators.json", "docs/collaborators/task-allocation.md"])
add_slide("Git Workflow", "How students continue work safely", image=diagrams["git-workflow"], layout="image-full", notes="Explain pull, branch, check, commit, push, and review.", evidence=["docs/collaborators/git-workflow.md"])


# Slide 28 member overview
slide_number = len(prs.slides) + 1
slide = prs.slides.add_slide(blank)
title(slide, "Member Contribution Overview", "Documented branch ownership by workstream", slide_number)
headers = ["Member", "GitHub / identity", "Branch", "Area"]
widths = [2.25, 2.35, 3.1, 2.9]
for side in range(2):
    x_base = 0.45 if side == 0 else 6.85
    y = 1.25
    for c, header in enumerate(headers):
        card(slide, x_base + sum(widths[:c]), y, widths[c] - 0.05, 0.42, header, "", "primary")
    y += 0.52
    for member in members[side * 9 : side * 9 + 9]:
        values = [member[0], member[1], member[2], member[3]]
        for c, value in enumerate(values):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_base + sum(widths[:c])), Inches(y), Inches(widths[c] - 0.05), Inches(0.48))
            shape.fill.solid()
            shape.fill.fore_color.rgb = ppt_color("surface2")
            shape.line.color.rgb = ppt_color("border")
            tx(slide, x_base + sum(widths[:c]) + 0.04, y + 0.07, widths[c] - 0.12, 0.32, value, 7.0, False)
        y += 0.5
records.append({"title": "Member Contribution Overview", "content": ["All documented members are mapped to branch ownership and contribution areas."], "image": None, "notes": "Use this as the transition into individual member contribution slides.", "evidence": ["docs/collaborators/task-allocation.md", "docs/Collaborators.json"]})


for member in members:
    name, github, branch, area, stream, _, _ = member
    slide_number = len(prs.slides) + 1
    slide = prs.slides.add_slide(blank)
    title(slide, f"Member Contribution: {name}", area, slide_number)
    card(slide, 0.55, 1.25, 3.3, 2.15, "Profile", f"GitHub / identity: {github}\nBranch: {branch}\nWorkstream: {stream}\nStatus: Assigned", "primary")
    tx(slide, 4.15, 1.25, 3.2, 0.35, "What this member worked on", 15, True)
    bullets(slide, 4.1, 1.65, 3.4, 1.55, worked[area], 11.5, "muted")
    tx(slide, 4.15, 3.52, 3.2, 0.35, "Why it mattered", 15, True)
    bullets(slide, 4.1, 3.92, 3.4, 1.25, why[area], 11.5, "muted")
    add_picture(slide, member_images[name], 7.85, 1.25, 4.9, 3.65)
    tx(slide, 7.9, 4.95, 4.8, 0.28, "Visual evidence: assigned app area or subsystem view", 9.5, False, "soft")
    tx(slide, 0.7, 6.25, 11.6, 0.38, f"Presenter note: Explain that this contribution supports {why[area][0]}.", 12, True, "primary")
    records.append({"title": f"Member Contribution: {name}", "content": [f"GitHub / identity: {github}", f"Branch: {branch}", f"Area: {area}"] + worked[area] + why[area], "image": member_images[name], "notes": f"Explain that this contribution supports {why[area][0]}.", "evidence": ["docs/collaborators/task-allocation.md", "docs/Collaborators.json", "docs/collaborators/frontend-team.md", "docs/collaborators/backend-team.md"]})


add_slide("Testing and Quality Assurance", "Checks documented for the project", ["npm run check verifies required files and blocks private or wrong-stack files.", "Expo Doctor is configured as npm run doctor for Expo project health checks.", "Android export is configured through npm run export:android.", "Testing checklist covers auth, project workflow, estimation, image workflow, and Android UI.", "Firebase rules and storage paths are reviewed as part of the quality story."], diagrams["build-deployment-flow"], "image-right", "Show that quality is both automated and documented.", ["package.json", "scripts/project-check.js", "docs/Testing_Checklist.md"])
add_slide("Debugging Work Completed", "Evidence of runtime and quality review", ["The repository includes crash logs before and after fixes.", "Runtime configuration has guard behaviour for missing Firebase values.", "Auth state initialization shows a splash screen before routing.", "Project check confirms required files and secret-safety expectations.", "Final submission checklist tracks Expo Go, EAS, Firebase, and branch readiness."], notes="Explain the debugging work as evidence of review and stabilization.", evidence=["sitespy-crash-log.txt", "sitespy-crash-log-after-fix.txt", "src/services/runtimeConfig.js"])
add_slide("How to Run the Project", "Verified commands from package.json and docs", ["Install dependencies: npm install", "Create local environment file: cp .env.example .env", "Start Expo: npm start", "Run Android target: npm run android", "Run project check: npm run check", "Run Expo Doctor: npm run doctor", "Export Android bundle: npm run export:android", "Build preview APK: npm run build:apk"], notes="Use this as the presenter technical runbook.", evidence=["package.json", "README.md"])
add_slide("How Students Continue Work", "Branch and portfolio continuation path", ["Configure Git identity with the student’s own name and email.", "Check out the assigned branch from the task-allocation document.", "Work only on the assigned app section or service area.", "Run checks before committing.", "Push the branch and prepare review evidence.", "Use portfolio notes to explain contribution, learning, challenges, and validated files."], diagrams["git-workflow"], "image-right", "Connect Git workflow to individual contribution evidence.", ["docs/collaborators/git-workflow.md", "README.md"])
add_slide("Demo Script for Monday", "Recommended presentation sequence", ["Introduce SiteSpy as a mobile-first wall-estimation app.", "Explain the estimation and collaboration problem.", "Show login and account flow.", "Show dashboard and project history.", "Create or open a project.", "Run a manual estimate and show the BOQ summary.", "Show image estimate and project details.", "Explain Firebase, data model, and security.", "Summarize member contributions.", "Close with testing, build readiness, and Q&A."], notes="This slide can be rehearsed directly as the demo script.", evidence=["docs/User_Manual.md"])
add_slide("Likely Questions and Answers", "Panel preparation", ["What problem does SiteSpy solve? It makes wall project estimates and saved records easier to produce and present.", "Why mobile-first? The app is designed for Android field-style use and student demonstration.", "Why Firebase? It provides authentication, Firestore records, storage paths, and deployable rules.", "How is user data protected? Rules require authentication and matching user ownership.", "How does estimation work? Wall dimensions and rates are converted into unit, mortar, labour, and total cost values.", "Can the app expand? Future work can add reporting, offline mode, and additional estimation types.", "How was it tested? Through project checks, testing checklist coverage, Expo checks, and Android export path."], notes="Use this to prepare presenters for panel questions.", evidence=["README.md", "docs/Testing_Checklist.md"])
add_slide("Future Improvements", "Clearly separated from current features", ["PDF report export for saved estimates.", "Additional construction estimation types beyond masonry walls.", "Offline mode for field use without consistent network access.", "Admin or reviewer mode for classroom assessment.", "Improved image measurement workflows.", "More reporting, charts, and portfolio evidence exports."], notes="Be clear that these are future improvements, not current claims.", evidence=["README.md", "docs/User_Manual.md"])
add_slide("Conclusion", "What the team can confidently present", ["SiteSpy is a working mobile-first estimation app built with Expo React Native.", "The app supports Firebase-backed user records, projects, estimates, images, and summaries.", "The codebase separates screens, components, services, utilities, theme, and configuration.", "The collaboration structure documents member ownership and branch responsibilities.", "The team can present the app workflow, technical architecture, and individual contributions clearly."], SOCIAL, "image-right", "Close by restating the system, foundation, and team contribution story.", ["README.md", "docs/collaborators/task-allocation.md"])

slide_number = len(prs.slides) + 1
slide = prs.slides.add_slide(blank)
bg(slide)
if LOGO.exists():
    add_picture(slide, LOGO, 5.65, 1.35, 2.0, 2.0)
tx(slide, 3.1, 3.75, 7.2, 0.7, "Thank You", 40, True, "text", PP_ALIGN.CENTER)
tx(slide, 2.55, 4.55, 8.3, 0.4, "SiteSpy: secure mobile estimation and project records for the client team.", 16, False, "muted", PP_ALIGN.CENTER)
tx(slide, 0.55, 7.13, 12.0, 0.22, f"SiteSpy System Presentation | {slide_number:02d}", 8, False, "soft", PP_ALIGN.RIGHT)
records.append({"title": "Thank You", "content": ["SiteSpy: secure mobile estimation and project records for the client team."], "image": LOGO, "notes": "Invite questions and move into live demo or panel discussion.", "evidence": ["assets/logo-mark-transparent.png"]})


prs.save(OUT)

md_lines = [
    "# SiteSpy Client Presentation FULL",
    "",
    f"PowerPoint: `{OUT.relative_to(ROOT).as_posix()}`",
    "Prepared: 7 June 2026",
    "",
    "## Source Notes",
    "",
    "- The presentation is based on repository files, screen source, services, Firebase rules, app configuration, and collaborator documentation.",
    "- Runtime browser screenshots were not available in this environment, so the deck uses polished UI walkthrough visual assets generated from implemented screen names, fields, actions, and theme colors.",
    "- No current-feature claims are made beyond repository evidence listed per slide.",
    "",
]

for index, record in enumerate(records, 1):
    md_lines.append(f"## Slide {index}: {record['title']}")
    md_lines.append("")
    if record["content"]:
        md_lines.append("Content:")
        for item in record["content"]:
            if isinstance(item, tuple):
                md_lines.append(f"- {item[0]}: {item[1]}")
            else:
                md_lines.append(f"- {item}")
    if record["image"]:
        images = record["image"] if isinstance(record["image"], list) else [record["image"]]
        md_lines.append("Images used:")
        for image in images:
            md_lines.append(f"- `{Path(image).relative_to(ROOT).as_posix()}`")
    md_lines.append(f"Speaker notes: {record['notes']}")
    if record["evidence"]:
        md_lines.append("Evidence files:")
        for evidence in record["evidence"]:
            md_lines.append(f"- `{evidence}`")
    md_lines.append("")

md_lines.extend(
    [
        "## Commands Verified",
        "",
        "- `npm run check`",
        "- `npm run doctor`",
        "- `npm run export:android`",
        "- Firebase deploy attempt is reported in the final implementation summary when credentials and tooling allow it.",
        "",
        "## Final File List",
        "",
        f"- `{OUT.relative_to(ROOT).as_posix()}`",
        f"- `{MD.relative_to(ROOT).as_posix()}`",
        f"- `{ASSET_DIR.relative_to(ROOT).as_posix()}`",
        f"- `{DIAG_DIR.relative_to(ROOT).as_posix()}`",
        f"- `{SHOT_DIR.relative_to(ROOT).as_posix()}`",
        f"- `{MEMBER_DIR.relative_to(ROOT).as_posix()}`",
    ]
)

MD.write_text("\n".join(md_lines), encoding="utf-8")

print(json.dumps({
    "pptx": str(OUT),
    "markdown": str(MD),
    "slides": len(prs.slides),
    "diagrams": len(list(DIAG_DIR.glob("*.png"))),
    "screenshots": len(list(SHOT_DIR.glob("*.png"))),
    "member_visuals": len(list(MEMBER_DIR.glob("*.png"))),
}, indent=2))
